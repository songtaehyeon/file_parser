import os
import zipfile
import shutil
import xml.etree.ElementTree as ET
from datetime import datetime
import pytz
from docx_parser import DocumentParser
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pandas as pd
import openpyxl
from tkinter import filedialog


#Open XML 구조의 경우 zip파일 형태와 같기 때문에 zip으로 압축해제 가능
def change_zip(files):
    origin_extension = (files.split('.'))[1]
    change_extension = "zip"
    zip_name = (files.split('.'))[0]+"."+change_extension
    os.rename(files, zip_name)
    output = "./sample"
    format = "zip"
    shutil.unpack_archive(zip_name, output, format) #shutil 모듈을 이용해서 아카이브 풀기
    os.rename(zip_name, (zip_name.split('.'))[0]+'.'+origin_extension)

#파일 내부 구조를 통한 파일 확인
def check_extension(file_extension):
    fin_check = ['word', 'ppt', 'xl']
    file_root = './sample'
    os.chdir(file_root)
    dir_names = os.listdir('./')
    for i in range(len(fin_check)):
        if fin_check[i] in dir_names:
            if i == 0:
                print(f"이 파일의 확장자는 {file_extension}이고, 분석 결과 워드 파일로 나왔습니다")
                return os.getcwd()
            elif i == 1:
                print(f"이 파일의 확장자는 {file_extension}이고, 분석 결과 ppt 파일로 나왔습니다")
                return os.getcwd()
            elif i == 2:
                print(f"이 파일의 확장자는 {file_extension}이고, 분석 결과 엑셀 파일로 나왔습니다")
                return os.getcwd()

#한국 시간으로 맞추기
def convert_to_korean_time(utc_time_str):
    utc_time = datetime.strptime(utc_time_str, "%Y-%m-%dT%H:%M:%SZ")
    utc_time = pytz.utc.localize(utc_time)
    korean_tz = pytz.timezone('Asia/Seoul')
    korean_time = utc_time.astimezone(korean_tz)
    return korean_time.strftime("%Y-%m-%d %H:%M:%S KST")

#data 추출 
def file_info1():
    tree = ET.parse('./sample/docProps/core.xml')
    root = tree.getroot()

    creator = root.find('.//{http://purl.org/dc/elements/1.1/}creator').text
    lastModifiedBy = root.find('.//{http://schemas.openxmlformats.org/package/2006/metadata/core-properties}lastModifiedBy').text
    revision = root.find('.//{http://schemas.openxmlformats.org/package/2006/metadata/core-properties}revision').text
    created_value = root.find(".//dcterms:created", namespaces={'dcterms': 'http://purl.org/dc/terms/'}).text
    modified_value = root.find(".//dcterms:modified", namespaces={'dcterms': 'http://purl.org/dc/terms/'}).text

    print("제작자:", creator)
    print("마지막으로 수정한 사람:", lastModifiedBy)
    print("수정 횟수:", revision)
    print("처음 파일 생성 시간:", convert_to_korean_time(created_value))
    print("마지막 파일 수정 시간:", convert_to_korean_time(modified_value))
    print("\n")
    print("==파일내용==")

def file_info2():
    tree = ET.parse('./sample/docProps/core.xml')
    root = tree.getroot()

    creator = root.find('.//{http://purl.org/dc/elements/1.1/}creator').text
    lastModifiedBy = root.find('.//{http://schemas.openxmlformats.org/package/2006/metadata/core-properties}lastModifiedBy').text
    created_value = root.find(".//dcterms:created", namespaces={'dcterms': 'http://purl.org/dc/terms/'}).text
    modified_value = root.find(".//dcterms:modified", namespaces={'dcterms': 'http://purl.org/dc/terms/'}).text

    print("제작자:", creator)
    print("마지막으로 수정한 사람:", lastModifiedBy)
    print("처음 파일 생성 시간:", convert_to_korean_time(created_value))
    print("마지막 파일 수정 시간:", convert_to_korean_time(modified_value))
    print("\n")
    print("==파일내용==")

# 파일에서 XML 내용 읽기
def extract_text_from_xml_file(file_path):
    with open(file_path, 'r', encoding='utf-8') as file:
        xml_string = file.read()

    # 태그를 찾아 텍스트 추출 -> 이 태그는 데이터가 들어가 있음
    start_tag = '<w:t'
    end_tag = '</w:t>'
    texts = []

    while start_tag in xml_string:
        # 시작 태그의 시작 및 끝 위치 찾기
        start_index = xml_string.find(start_tag)
        start_close_index = xml_string.find('>', start_index)
        # 종료 태그 위치 찾기
        end_index = xml_string.find(end_tag, start_close_index)

        # 텍스트 추출 및 저장
        text = xml_string[start_close_index+1:end_index].strip()
        texts.append(text)

        # 다음 태그 검색을 위해 문자열 자르기
        xml_string = xml_string[end_index+len(end_tag):]
    return texts

def docx_docm_info(files):
    file_info1()
    change_zip(files)
    if os.path.exists('./sample/word/vbaProject.bin'):
        print("※주의 : 매크로가 포함된 파일입니다.")
        file_info1()
        texts = extract_text_from_xml_file('./sample/word/document.xml') # 이 방식을 택한 이유는 documentparser가 docm 을 지원 안하기 때문
        for text in texts:
            print(text)

    else:
        file_info1()
        doc = DocumentParser(files)
        for _type, item in doc.parse():
            print(_type, item)

def pptx_info(files):
    file_info1()
    parsed = Presentation(files)
    for slide in parsed.slides:
        for shape in slide.shapes:
            print(shape.shape_type)
            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                print("텍스트 박스 안의 내용:", shape.text_frame.text)
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                row_count = len(shape.table.rows)
                col_count = len(shape.table.columns)
                table_data = []
                for _r in range(0, row_count):
                    row = []
                    for _c in range(0, col_count):
                        cell = shape.table.cell(_r, _c)
                        row.append(cell.text)
                    table_data.append(row)
                print(table_data)
                df_temp = pd.DataFrame(columns=table_data[0], data=table_data[1:])
                print(df_temp)

def xlsx_info(files):
    file_info2()
    workbook = openpyxl.load_workbook(files)
    sheet = workbook.active

    for row in sheet.iter_rows(values_only=True):
        print('\t'.join(str(cell_value) for cell_value in row))

def main():
    list_file = []
    files = filedialog.askopenfilename(initialdir="./",\
                 title = "파일을 선택 해 주세요",\
                    filetypes = (("all","*"),("*.xlsx","*xlsx"),("*.pptx","*pptx"),("*.docx","*docx"), ("*.docm","*docm")))

    file_extension = (files.split('.'))[1]

    if file_extension in office_type:
        change_zip(files)
        check_extension(file_extension)
        print("\n")
        os.chdir('..')
        if file_extension in ['docx', 'docm']:
            docx_docm_info(files)
        elif file_extension == 'pptx':
            pptx_info(files)
        elif file_extension == 'xlsx':
            xlsx_info(files)

office_type = ["docx", "xlsx", "pptx", "docm"]
main()
